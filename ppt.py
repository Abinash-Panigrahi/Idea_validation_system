"""
ppt.py — Pitch Deck Generator for ThynxAI
Opens template_1.pptx and injects Gemini data into placeholders.
No API calls here — just data injection.
"""

import io
import os
import re
from pptx import Presentation


# ─── Template Path ────────────────────────────────────────────────────────────

TEMPLATE_PATH = os.path.join(
    os.path.dirname(__file__),
    "ppt_templates",
    "template_1.pptx"
)


# ─── Helper: Replace text in a shape ─────────────────────────────────────────

def replace_in_shape(shape, data: dict):
    if not shape.has_text_frame:
        return

    for para in shape.text_frame.paragraphs:
        # Join all runs into one full text first
        full_text = "".join([run.text for run in para.runs])
        
        # Check if any placeholder exists in full text
        if "{{" not in full_text:
            continue
        
        # Replace all placeholders in full text
        for key, value in data.items():
            full_text = full_text.replace(f"{{{{{key}}}}}", str(value))
        
        # Write back to first run, clear the rest
        if para.runs:
            para.runs[0].text = full_text
            for run in para.runs[1:]:
                run.text = ""


def replace_in_table(shape, data: dict):
    if not shape.has_table:
        return

    for row in shape.table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                full_text = "".join([run.text for run in para.runs])
                
                if "{{" not in full_text:
                    continue
                
                for key, value in data.items():
                    full_text = full_text.replace(f"{{{{{key}}}}}", str(value))
                
                if para.runs:
                    para.runs[0].text = full_text
                    for run in para.runs[1:]:
                        run.text = ""


# ─── Main Generator ───────────────────────────────────────────────────────────

def generate_ppt(slide_data: dict) -> bytes:
    """
    Opens template_1.pptx, injects data into all {{placeholders}}.
    Returns .pptx file as bytes for Streamlit download.
    """
    if not os.path.exists(TEMPLATE_PATH):
        raise FileNotFoundError(
            f"Template not found at {TEMPLATE_PATH}. "
            "Make sure template_1.pptx is inside the ppt_templates folder."
        )

    prs = Presentation(TEMPLATE_PATH)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_in_shape(shape, slide_data)
            elif shape.has_table:
                replace_in_table(shape, slide_data)

    # Save to bytes for Streamlit in-memory download
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()